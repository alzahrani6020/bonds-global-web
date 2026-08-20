#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Create funding-focused assets for investor pitch.
Generates charts, slide content, and a PowerPoint deck emphasizing investment and returns.
"""

import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib import font_manager as fm
import numpy as np
import arabic_reshaper
from bidi.algorithm import get_display
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Bonds brand colors
BONDS_NAVY = '#0A0F1A'
BONDS_GOLD = '#D4A853'
BONDS_GOLD_BRIGHT = '#F0C96A'
BONDS_TEXT = '#E8ECF4'
BONDS_TEXT_SECONDARY = '#94A3B8'
BONDS_GREEN = '#22C55E'
BONDS_RED = '#EF4444'
BONDS_BLUE = '#3B82F6'
BONDS_PURPLE = '#A855F7'
BONDS_TEAL = '#14B8A6'

OUTPUT_DIR = 'canva-ready-funding'
CHARTS_DIR = os.path.join(OUTPUT_DIR, 'charts')


def ar(text):
    """Reshape and reorder Arabic text for matplotlib/PPTX rendering."""
    if not isinstance(text, str):
        text = str(text)
    reshaped = arabic_reshaper.reshape(text)
    return get_display(reshaped)


def setup_arabic_font():
    possible_fonts = ['Arial', 'Segoe UI', 'Tahoma', 'Calibri', 'Noto Sans Arabic']
    chosen_font = None
    for font_name in possible_fonts:
        try:
            font_path = fm.findfont(fm.FontProperties(family=font_name))
            if font_path and font_path.lower().endswith(('.ttf', '.ttc', '.otf')):
                chosen_font = font_name
                break
        except Exception:
            continue
    if chosen_font:
        plt.rcParams['font.family'] = chosen_font
    plt.rcParams['axes.unicode_minus'] = False


setup_arabic_font()
os.makedirs(CHARTS_DIR, exist_ok=True)


def save_chart(fig, name):
    fig.savefig(
        os.path.join(CHARTS_DIR, name),
        dpi=300,
        bbox_inches='tight',
        facecolor=BONDS_NAVY,
        edgecolor='none'
    )
    plt.close(fig)
    print(f"Saved chart: {name}")


def style_axis(ax):
    ax.set_facecolor(BONDS_NAVY)
    ax.tick_params(axis='both', labelsize=11)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color(BONDS_TEXT_SECONDARY)
    ax.spines['bottom'].set_color(BONDS_TEXT_SECONDARY)


def create_capital_structure():
    fig, ax = plt.subplots(figsize=(12, 7), facecolor=BONDS_NAVY)
    style_axis(ax)

    labels = [ar('رأس المال\nالمصرح'), ar('الاحتياطي\nالعام'), ar('رأس المال\nالمستثمر'), ar('صافي\nالتمويل')]
    values = [200, -50, 226, 176]
    colors = [BONDS_GOLD, BONDS_RED, BONDS_GOLD_BRIGHT, BONDS_GREEN]
    cumulative = np.cumsum([0] + values[:-1])

    bars = ax.bar(labels, values, bottom=cumulative, color=colors, edgecolor='white', linewidth=1.5, width=0.6)
    for bar, val, cum in zip(bars, values, cumulative):
        ax.text(bar.get_x() + bar.get_width()/2, cum + val/2,
                ar(f'{abs(val)} مليون ر.س'), ha='center', va='center',
                fontsize=13, fontweight='bold', color=BONDS_NAVY)

    ax.set_title(ar('هيكل رأس المال المستثمر'), fontsize=22, fontweight='bold', color=BONDS_GOLD, pad=20)
    ax.set_ylabel(ar('مليون ريال سعودي'), fontsize=14, color=BONDS_TEXT_SECONDARY)
    ax.tick_params(axis='x', labelsize=13)
    ax.set_ylim(0, 280)
    ax.axhline(y=0, color=BONDS_TEXT_SECONDARY, linestyle='-', linewidth=0.5)
    save_chart(fig, 'chart-capital-structure.png')


def create_use_of_funds():
    fig, ax = plt.subplots(figsize=(12, 8), facecolor=BONDS_NAVY)
    ax.set_facecolor(BONDS_NAVY)

    labels = [ar('شراء وتشطيب\nالأصول'), ar('احتياطي\nتشغيلي'), ar('التسويق\nوالمبيعات'),
              ar('استشارات\nوقانون'), ar('تقنية\nوأنظمة'), ar('إدارة\nوتأسيس')]
    sizes = [55, 20, 10, 5, 5, 5]
    colors = [BONDS_GOLD, BONDS_GOLD_BRIGHT, BONDS_BLUE, BONDS_TEAL, BONDS_PURPLE, BONDS_TEXT_SECONDARY]

    wedges, texts, autotexts = ax.pie(
        sizes, labels=labels, colors=colors, autopct='%1.0f%%',
        startangle=90, pctdistance=0.75,
        wedgeprops=dict(width=0.5, edgecolor=BONDS_NAVY, linewidth=2)
    )
    for text in texts:
        text.set_fontsize(11)
        text.set_color(BONDS_TEXT)
    for autotext in autotexts:
        autotext.set_fontsize(11)
        autotext.set_fontweight('bold')
        autotext.set_color(BONDS_NAVY)

    ax.set_title(ar('استخدامات رأس المال (226 مليون ر.س)'), fontsize=20, fontweight='bold', color=BONDS_GOLD, pad=20)
    save_chart(fig, 'chart-use-of-funds.png')


def create_irr_projection():
    fig, ax = plt.subplots(figsize=(13, 7), facecolor=BONDS_NAVY)
    style_axis(ax)

    years = np.arange(1, 11)
    irr_pessimistic = [2, 5, 8, 10, 12, 13, 14, 14.5, 15, 16]
    irr_expected = [5, 10, 14, 16, 18, 19, 19.5, 20, 20, 20]
    irr_optimistic = [8, 15, 20, 23, 25, 26, 27, 28, 29, 30]

    ax.fill_between(years, irr_pessimistic, irr_optimistic, alpha=0.15, color=BONDS_GOLD)
    ax.plot(years, irr_expected, color=BONDS_GOLD, linewidth=3, marker='o', markersize=8, label=ar('السيناريو المتوقع'))
    ax.plot(years, irr_pessimistic, color=BONDS_RED, linewidth=2, linestyle='--', label=ar('السيناريو المتشائم'))
    ax.plot(years, irr_optimistic, color=BONDS_GREEN, linewidth=2, linestyle='--', label=ar('السيناريو المتفائل'))

    ax.axhline(y=19, color=BONDS_GOLD_BRIGHT, linestyle=':', linewidth=1.5, alpha=0.8)
    ax.text(9.5, 19.8, ar('هدف IRR 19%'), color=BONDS_GOLD_BRIGHT, fontsize=11, ha='right')

    ax.set_xlabel(ar('السنة'), fontsize=14, color=BONDS_TEXT_SECONDARY)
    ax.set_ylabel(ar('العائد الداخلي (%)'), fontsize=14, color=BONDS_TEXT_SECONDARY)
    ax.set_title(ar('توقعات العائد الداخلي على 10 سنوات'), fontsize=20, fontweight='bold', color=BONDS_GOLD, pad=20)
    ax.legend(loc='upper left', fontsize=11, facecolor=BONDS_NAVY, edgecolor=BONDS_GOLD)
    ax.set_xticks(years)
    ax.set_ylim(0, 35)
    save_chart(fig, 'chart-irr-projection.png')


def create_exit_strategy():
    fig, ax = plt.subplots(figsize=(13, 6), facecolor=BONDS_NAVY)
    style_axis(ax)

    phases = [ar('السنة 1-2'), ar('السنة 3-4'), ar('السنة 5-7'), ar('السنة 8-10')]
    returns = [1.0, 1.5, 2.5, 4.0]
    colors = [BONDS_TEXT_SECONDARY, BONDS_BLUE, BONDS_GOLD, BONDS_GREEN]

    bars = ax.barh(phases, returns, color=colors, edgecolor='white', linewidth=1.5, height=0.6)
    for bar, ret in zip(bars, returns):
        ax.text(bar.get_width() + 0.1, bar.get_y() + bar.get_height()/2,
                f'{ret}x', va='center', fontsize=16, fontweight='bold', color=BONDS_GOLD_BRIGHT)

    ax.set_xlabel(ar('مضاعفة رأس المال'), fontsize=14, color=BONDS_TEXT_SECONDARY)
    ax.set_title(ar('خطة الخروج الاستراتيجي'), fontsize=20, fontweight='bold', color=BONDS_GOLD, pad=20)
    ax.set_xlim(0, 5)
    save_chart(fig, 'chart-exit-strategy.png')


def create_investment_comparison():
    fig, ax = plt.subplots(figsize=(12, 7), facecolor=BONDS_NAVY)
    style_axis(ax)

    models = [ar('العقار\nالتقليدي'), ar('الصندوق\nالعقاري'), ar('شركة\nبوندز')]
    irr = [8, 12, 20]
    exit_years = [15, 10, 7]
    min_invest = [50, 25, 5]

    x = np.arange(len(models))
    width = 0.25

    ax.bar(x - width, irr, width, label=ar('IRR %'), color=BONDS_GOLD)
    ax.bar(x, [y/2 for y in exit_years], width, label=ar('سنوات الخروج ÷2'), color=BONDS_BLUE)
    ax.bar(x + width, [m*2 for m in min_invest], width, label=ar('الحد الأدنى ×2 مليون'), color=BONDS_TEAL)

    ax.set_ylabel(ar('القيمة النسبية'), fontsize=14, color=BONDS_TEXT_SECONDARY)
    ax.set_title(ar('مقارنة فرص الاستثمار'), fontsize=20, fontweight='bold', color=BONDS_GOLD, pad=20)
    ax.set_xticks(x)
    ax.set_xticklabels(models, fontsize=12)
    ax.legend(loc='upper right', fontsize=11, facecolor=BONDS_NAVY, edgecolor=BONDS_GOLD)
    save_chart(fig, 'chart-investment-comparison.png')


def create_cash_flow():
    fig, ax = plt.subplots(figsize=(13, 7), facecolor=BONDS_NAVY)
    style_axis(ax)

    years = np.arange(0, 11)
    cash_flow = [-50, -30, -20, 10, 35, 60, 85, 110, 140, 175, 210]
    cumulative = np.cumsum(cash_flow)
    positive = [max(0, v) for v in cash_flow]
    negative = [min(0, v) for v in cash_flow]

    ax.bar(years, positive, color=BONDS_GREEN, edgecolor=BONDS_NAVY, linewidth=1.5, label=ar('التدفق الموجب'))
    ax.bar(years, negative, color=BONDS_RED, edgecolor=BONDS_NAVY, linewidth=1.5, label=ar('التدفق السالب'))
    ax.plot(years, cumulative, color=BONDS_GOLD, linewidth=3, marker='o', markersize=6, label=ar('الرصيد التراكمي'))
    ax.axhline(y=0, color=BONDS_TEXT_SECONDARY, linestyle='-', linewidth=0.8)

    ax.set_xlabel(ar('السنة'), fontsize=14, color=BONDS_TEXT_SECONDARY)
    ax.set_ylabel(ar('مليون ريال سعودي'), fontsize=14, color=BONDS_TEXT_SECONDARY)
    ax.set_title(ar('توقع التدفق النقدي - السيناريو المتوقع'), fontsize=20, fontweight='bold', color=BONDS_GOLD, pad=20)
    ax.legend(loc='upper left', fontsize=11, facecolor=BONDS_NAVY, edgecolor=BONDS_GOLD)
    ax.set_xticks(years)
    save_chart(fig, 'chart-cash-flow-projection.png')


# =============================================================================
# PowerPoint generation
# =============================================================================
def hex_to_rgb(hex_color):
    hex_color = hex_color.lstrip('#')
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def add_background(slide, prs):
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = hex_to_rgb(BONDS_NAVY)
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_text_box(slide, left, top, width, height, text, font_size, font_color, bold=False, align=PP_ALIGN.RIGHT):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = ar(text)
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.color.rgb = hex_to_rgb(font_color)
    run.font.bold = bold
    run.font.name = 'Arial'
    return box


def add_bullet_box(slide, left, top, width, height, bullets, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = ar('• ' + bullet)
        p.alignment = PP_ALIGN.RIGHT
        p.level = 0
        run = p.runs[0]
        run.font.size = Pt(font_size)
        run.font.color.rgb = hex_to_rgb(BONDS_TEXT)
        run.font.name = 'Arial'
    return box


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    logo_path = 'bonds-logo.png'
    qr_path = 'bonds-qr-code.png'
    has_logo = os.path.exists(logo_path)
    has_qr = os.path.exists(qr_path)

    slides_data = [
        {
            'title': 'فرصة استثمارية',
            'subtitle': 'إحياء الأصول العقارية غير المكتملة',
            'bullets': [
                'عائد داخلي مستهدف 19% - 20%',
                'مضاعفة رأس المال 2.5x - 4x',
                'تمويل ذاتي 100% · لا قروض',
                'الحد الأدنى للمساهمة: 5 ملايين ريال'
            ]
        },
        {
            'title': 'لماذا التمويل الآن؟',
            'bullets': [
                'نقص في العرض ووفرة في الطلب في السوق السعودي',
                'رؤية 2030 تدعم القطاع العقاري بشكل غير مسبوق',
                'أكثر من 7,000 عقار غير مكتمل في الرياض وجدة',
                'نموذج المحاصة يقلل المخاطر ويزيد العوائد',
                'لا حاجة لشراء أراضٍ — نبدأ من الأصول القائمة'
            ]
        },
        {
            'title': 'حجم الفرصة',
            'bullets': [
                'حجم التمويل العقاري: 730 مليار ريال',
                'نمو الصفقات العقارية: +14.2% سنوياً',
                'عائد إيجاري في الرياض: 8.89%',
                'عائد إيجاري في جدة: 7.89%',
                'سوق ضخم يبحث عن حلول مبتكرة'
            ]
        },
        {
            'title': 'نموذج الشراكة (المحاصة)',
            'bullets': [
                'المالك يقدم العقار غير المكتمل',
                'بوندز تقدم التمويل والإدارة والتشطيب',
                'توزيع الأرباح حسب نسبة المساهمة',
                'عقود واضحة مع تحكيم ملزم',
                'كفالات بنكية وتأمين شامل'
            ]
        },
        {
            'title': 'هيكل رأس المال المستثمر',
            'bullets': [
                'رأس المال المصرح: 200 مليون ريال',
                'الاستثمار المستهدف للمرحلة الأولى: 226 مليون ريال',
                'التمويل: 100% ذاتي من المساهمين والمستثمرين',
                'لا قروض مصرفية · لا صكوك · لا تمويل ربوي',
                'الحد الأدنى للمساهمة: 5 ملايين ريال'
            ]
        },
        {
            'title': 'استخدامات رأس المال',
            'bullets': [
                'شراء وتشطيب الأصول: 55% (124 مليون ر.س)',
                'احتياطي تشغيلي: 20% (45 مليون ر.س)',
                'التسويق والمبيعات: 10% (23 مليون ر.س)',
                'الاستشارات والقانون: 5% (11 مليون ر.س)',
                'التقنية والأنظمة: 5% (11 مليون ر.س)',
                'الإدارة والتأسيس: 5% (11 مليون ر.س)'
            ]
        },
        {
            'title': 'توقعات العائد — IRR',
            'bullets': [
                'السيناريو المتشائم: IRR 16% · مضاعفة 2x',
                'السيناريو المتوقع: IRR 19% - 20% · مضاعفة 2.5x - 3x',
                'السيناريو المتفائل: IRR 27% - 30% · مضاعفة 4x',
                'أفق الاستثمار: 7 - 10 سنوات',
                'نقطة التعادل: السنة 3 - 4'
            ]
        },
        {
            'title': 'توقع التدفق النقدي',
            'bullets': [
                'السنوات 1-2: استثمار وتشطيب',
                'السنة 3: التحول للإيجابية',
                'السنوات 4-7: نمو متسارع للتدفقات',
                'السنوات 8-10: استقرار وتوزيع أرباح'
            ]
        },
        {
            'title': 'خطة الخروج',
            'bullets': [
                'السنة 1-2: إثبات النموذج · مضاعفة 1x',
                'السنة 3-4: النمو والتكرار · مضاعفة 1.5x',
                'السنة 5-7: خروج جزئي لصندوق استثماري · مضاعفة 2.5x',
                'السنة 8-10: الطرح أو إعادة الشراء · مضاعفة 4x'
            ]
        },
        {
            'title': 'مقارنة فرص الاستثمار',
            'bullets': [
                'العقار التقليدي: IRR ~8% · دورة 15 سنة · 50 مليون',
                'الصندوق العقاري: IRR ~12% · دورة 10 سنوات · 25 مليون',
                'بوندز: IRR 19% - 20% · دورة 7 سنوات · من 5 ملايين'
            ]
        },
        {
            'title': 'إدارة المخاطر للمستثمر',
            'bullets': [
                'تأمين شامل على جميع الأصول',
                'فحص هندسي وقانوني قبل كل صفقة',
                'عقود محاصة واضحة مع تحكيم ملزم',
                'كفالات بنكية على المقاولين',
                'تنويع المحفظة على عدة أصول ومناطق',
                'فصل بين أموال المستثمرين والشركة'
            ]
        },
        {
            'title': 'الفريق المالي والإداري',
            'bullets': [
                'المستشار المالي: د. طلال بن حسن الزهراني',
                'خبرة 25+ عاماً في الأسواق المالية والعقارية',
                'فريق هندسي وقانوني متخصص',
                'نظام إداري ومالي شفاف',
                'تقارير دورية للمستثمرين'
            ]
        },
        {
            'title': 'الخطوات القادمة',
            'bullets': [
                'إغلاق جولة التمويل الأولى: 226 مليون ريال',
                'تأسيس الشركة والحصول على التراخيص',
                'إطلاق أول 3 مشاريع في الرياض وجدة',
                'بناء قاعدة بيانات الأصول غير المكتملة',
                'التوسع إلى 10 مدن سعودية خلال 5 سنوات'
            ]
        },
        {
            'title': 'الدعوة للاستثمار',
            'subtitle': 'كن شريكاً في بناء نموذج استثماري فريد',
            'bullets': [
                'الاستثمار يبدأ من 5 ملايين ريال',
                'عائد مستهدف 19% - 20%',
                'د. طلال بن حسن الزهراني',
                '+966 56 756 6616',
                'info@bonds-global.com',
                'www.bonds-global.com'
            ]
        },
    ]

    for idx, data in enumerate(slides_data):
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        add_background(slide, prs)

        # Logo top-right
        if has_logo:
            slide.shapes.add_picture(logo_path, Inches(11.2), Inches(0.3), width=Inches(1.6))

        # Slide number
        add_text_box(slide, 0.4, 0.3, 1, 0.4, str(idx + 1), 14, BONDS_TEXT_SECONDARY, align=PP_ALIGN.LEFT)

        # Title
        add_text_box(slide, 1, 0.5, 10, 1, data['title'], 36, BONDS_GOLD, bold=True)

        # Subtitle if exists
        top_offset = 1.5
        if 'subtitle' in data:
            add_text_box(slide, 1, 1.4, 10, 0.6, data['subtitle'], 22, BONDS_GOLD_BRIGHT)
            top_offset = 2.0

        # Bullets
        add_bullet_box(slide, 1, top_offset, 10, 4.5, data['bullets'], font_size=22)

        # QR code on last slide
        if idx == len(slides_data) - 1 and has_qr:
            slide.shapes.add_picture(qr_path, Inches(0.8), Inches(5.5), width=Inches(1.2))
            add_text_box(slide, 0.5, 6.7, 2, 0.4, 'info@bonds-global.com', 11, BONDS_TEXT_SECONDARY, align=PP_ALIGN.CENTER)

    pptx_path = os.path.join(OUTPUT_DIR, 'funding-focused-pitch-arabic.pptx')
    prs.save(pptx_path)
    print(f"Saved presentation: {pptx_path}")


if __name__ == '__main__':
    create_capital_structure()
    create_use_of_funds()
    create_irr_projection()
    create_exit_strategy()
    create_investment_comparison()
    create_cash_flow()
    create_presentation()

    slide_content = """# محتوى شرائح قسم التمويل المُعزّز

## الشريحة 1: الغلاف — فرصة استثمارية
- **العنوان:** فرصة استثمارية في إحياء الأصول العقارية غير المكتملة
- **العنوان الفرعي:** عائد داخلي مستهدف 19% - 20% · مضاعفة رأس المال 2.5x - 4x
- **الشعار:** شعار بوندز
- **الألوان:** خلفية داكنة #0A0F1A، نص ذهبي #D4A853

## الشريحة 2: لماذا التمويل الآن؟
- السوق السعودي يشهد نقصاً في العرض ووفرة في الطلب
- رؤية 2030 تدعم القطاع العقاري بشكل غير مسبوق
- آلاف الأصول غير المكتملة بحاجة لشريك مالي وإداري
- نموذج المحاصة يقلل المخاطر ويزيد العوائد
- لا حاجة لشراء أراضٍ — نبدأ من الأصول القائمة

## الشريحة 3: حجم الفرصة
- أكثر من 7,000 عقار غير مكتمل في الرياض وجدة فقط
- حجم التمويل العقاري: 730 مليار ريال
- نمو الصفقات العقارية: +14.2% سنوياً
- عائد إيجاري في الرياض: 8.89%
- عائد إيجاري في جدة: 7.89%

## الشريحة 4: نموذج الشراكة (المحاصة)
- المالك يقدم العقار غير المكتمل
- بوندز تقدم التمويل والإدارة والتشطيب
- توزيع الأرباح حسب نسبة المساهمة
- عقود واضحة مع تحكيم ملزم
- كفالات بنكية وتأمين شامل

## الشريحة 5: هيكل رأس المال المستثمر
- رأس المال المصرح: 200 مليون ريال
- الاستثمار المستهدف للمرحلة الأولى: 226 مليون ريال
- التمويل: 100% ذاتي من المساهمين المؤسسين والمستثمرين
- لا قروض مصرفية · لا صكوك · لا تمويل ربوي
- الحد الأدنى للمساهمة: 5 ملايين ريال

## الشريحة 6: استخدامات رأس المال
- شراء وتشطيب الأصول: 55% (124 مليون ر.س)
- احتياطي تشغيلي: 20% (45 مليون ر.س)
- التسويق والمبيعات: 10% (23 مليون ر.س)
- الاستشارات الهندسية والقانونية: 5% (11 مليون ر.س)
- التقنية والأنظمة: 5% (11 مليون ر.س)
- الإدارة والتأسيس: 5% (11 مليون ر.س)

## الشريحة 7: توقعات العائد — IRR
- السيناريو المتشائم: IRR 16% · مضاعفة 2x
- السيناريو المتوقع: IRR 19% - 20% · مضاعفة 2.5x - 3x
- السيناريو المتفائل: IRR 27% - 30% · مضاعفة 4x
- أفق الاستثمار: 7 - 10 سنوات
- نقطة التعادل: السنة 3 - 4

## الشريحة 8: توقع التدفق النقدي
- السنوات 1-2: استثمار وتشطيب (تدفق سلبي متوقع)
- السنة 3: بداية التحول للإيجابية
- السنوات 4-7: نمو التدفقات النقدية بوتيرة متسارعة
- السنوات 8-10: استقرار وتوزيع أرباح سنوي

## الشريحة 9: خطة الخروج
- **السنة 1-2:** إثبات النموذج · مضاعفة 1x
- **السنة 3-4:** النمو والتكرار · مضاعفة 1.5x
- **السنة 5-7:** خروج جزئي عبر البيع لصندوق استثماري · مضاعفة 2.5x
- **السنة 8-10:** الطرح أو إعادة الشراء · مضاعفة 4x

## الشريحة 10: مقارنة فرصة الاستثمار
- العقار التقليدي: IRR ~8% · دورة 15 سنة · استثمار 50 مليون
- الصندوق العقاري: IRR ~12% · دورة 10 سنوات · استثمار 25 مليون
- **بوندز: IRR 19% - 20% · دورة 7 سنوات · استثمار من 5 ملايين**

## الشريحة 11: إدارة المخاطر للمستثمر
- تأمين شامل على جميع الأصول
- فحص هندسي وقانوني قبل كل صفقة
- عقود محاصة واضحة مع تحكيم ملزم
- كفالات بنكية على المقاولين
- تنويع المحفظة على عدة أصول ومناطق
- فصل بين أموال المستثمرين وحسابات الشركة

## الشريحة 12: الفريق المالي والإداري
- المستشار المالي: د. طلال بن حسن الزهراني
- خبرة 25+ عاماً في الأسواق المالية والعقارية
- فريق هندسي وقانوني متخصص
- نظام إداري ومالي شفاف
- تقارير دورية للمستثمرين

## الشريحة 13: الخطوات القادمة
- إغلاق جولة التمويل الأولى: 226 مليون ريال
- تأسيس الشركة والحصول على التراخيص
- إطلاق أول 3 مشاريع في الرياض وجدة
- بناء قاعدة بيانات الأصول غير المكتملة
- التوسع إلى 10 مدن سعودية خلال 5 سنوات

## الشريحة 14: الدعوة للاستثمار
- **كن شريكاً في بناء نموذج استثماري فريد**
- الاستثمار يبدأ من 5 ملايين ريال
- عائد مستهدف 19% - 20%
- تواصل معنا:
  - د. طلال بن حسن الزهراني
  - +966 56 756 6616
  - info@bonds-global.com
  - www.bonds-global.com
"""

    with open(os.path.join(OUTPUT_DIR, 'slide-content-funding.md'), 'w', encoding='utf-8') as f:
        f.write(slide_content)

    print(f"\nSaved slide content: {os.path.join(OUTPUT_DIR, 'slide-content-funding.md')}")
    print(f"All funding assets created in: {OUTPUT_DIR}/")
