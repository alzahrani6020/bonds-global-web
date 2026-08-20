from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

BONDS_GOLD = RGBColor(212, 168, 83)
BONDS_DARK = RGBColor(10, 15, 26)
BONDS_GRAY = RGBColor(148, 163, 184)

def add_title_slide(prs, title, subtitle):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BONDS_DARK
    background.line.fill.background()
    
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(4.25), Inches(0.3), width=Inches(1.5))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BONDS_GOLD
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.2))
    tf = contact_box.text_frame
    tf.word_wrap = True
    lines = [
        ('Financial Advisor', 14, True, BONDS_GOLD),
        ('Dr. Talal bin Hassan Al-Zahrani', 16, True, RGBColor(255, 255, 255)),
        ('+966 56 756 6616', 12, False, BONDS_GRAY),
    ]
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BONDS_DARK
    header.line.fill.background()
    
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.15), width=Inches(0.8))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BONDS_GOLD
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT
    
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = '• ' + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = BONDS_DARK
    header.line.fill.background()
    
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.15), width=Inches(0.8))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8.5), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BONDS_GOLD
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT
    
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = '• ' + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
    
    right_box = slide.shapes.add_textbox(Inches(4.8), Inches(1.5), Inches(4.2), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.LEFT
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = '• ' + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
    
    return slide

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

add_title_slide(prs, 'Business Plan', 'Real Estate Asset Revitalization Company')

add_content_slide(prs, 'Executive Summary', [
    'Unique investment model for revitalizing incomplete real estate assets',
    'Partnership agreements with owners without purchasing land',
    'Self-funded by founding shareholders',
    'Compliance with Saudi Building Code and comprehensive insurance',
    'Expected IRR of 19% - 20% over 10 years'
])

add_two_column_slide(prs, 'Problem & Opportunity',
    'Problem', [
        'Thousands of incomplete buildings',
        'Owners lack liquidity',
        'Difficulty accessing reliable contractors',
        'Complex legal procedures'
    ],
    'Opportunity', [
        'Vision 2030 and real estate growth',
        'Demand exceeds supply in major cities',
        'High rental yields',
        'Flexible and profitable partnership model'
    ])

add_content_slide(prs, 'Solution: Business Model', [
    'Identify incomplete real estate assets in target locations',
    'Conduct engineering, legal, and financial assessment',
    'Negotiate and sign clear partnership contract',
    'Execute finishing works to highest quality standards',
    'Market and sell units at market prices',
    'Distribute profits between company and owner'
])

add_content_slide(prs, 'Real Estate Market Indicators', [
    'Rental yield in Riyadh: 8.89%',
    'Rental yield in Jeddah: 7.89%',
    'Real estate transaction growth (2025): +14.2%',
    'Real estate financing volume: SAR 730 billion',
    'Price-to-demand gap: 6%',
    'New construction premium: +25%'
])

add_two_column_slide(prs, 'Target Areas',
    'Riyadh', [
        'Premium North: Al-Aqeeq, Hittin, Al-Malqa',
        'Northeast: Al-Munsiyah, Al-Rimal',
        'Al-Yasmin, Al-Narjis',
        'High rental yields'
    ],
    'Jeddah', [
        'North Coast: Abhur, Al-Sawari',
        'Vibrant Center: Al-Hamra, Al-Naeem',
        'Al-Zahra, Al-Salama',
        'Continuous residential and commercial demand'
    ])

add_content_slide(prs, 'Partnership Models', [
    'Full Muhassah: Company handles finishing and marketing for profit share',
    'Partial Purchase: Company buys part of units after finishing',
    'Management Fee: Fixed fee for finishing and marketing',
    'Hybrid Model: Mix of Muhassah, purchase, and management fees'
])

add_content_slide(prs, 'Key Financial Indicators', [
    'Authorized Capital: SAR 200 million',
    'Total Target Investment: SAR 226 million',
    'Initial Investment: SAR 25 million',
    'Expected IRR: 19% - 20%',
    'Capital Return Multiple at Exit: 2.5x - 4x',
    'Payback Period: 5.5 years'
])

add_content_slide(prs, 'Revenue Sources', [
    'Muhassah revenues: 60%',
    'Resale and purchase profits: 25%',
    'Management and marketing fees: 10%',
    'Consulting and administrative services: 5%'
])

add_content_slide(prs, 'Implementation Timeline', [
    'Years 1-2: Proof of Concept (3-5 projects)',
    'Years 3-5: Growth and Replication (6-10 projects annually)',
    'Years 5-7: Strategic Exit',
    'Options: Sale of controlling stake, partial IPO, or buyback'
])

add_content_slide(prs, 'Risk Management', [
    'Comprehensive insurance (CAR, liability, hidden defects)',
    'Engineering and legal inspection before contracting',
    'Clear contracts with binding arbitration',
    'Bank guarantees from contractors',
    '10% risk reserve',
    'Geographic diversification and multiple target segments'
])

add_content_slide(prs, 'Management Team', [
    'CEO: Strategy and investor relations',
    'COO: Project management and quality',
    'CFO: Cash flows and reporting',
    'Commercial Director: Marketing and sales',
    'Legal Director: Contracts and compliance',
    'Hiring plan: 10-15 employees in Year 1'
])

add_content_slide(prs, 'Why Invest?', [
    'Unique model requiring no land purchase',
    'Large market of incomplete assets',
    'Self-funded without debt or sukuk',
    'Attractive returns with calculated risks',
    'Specialized management team and strong network',
    'Clear and flexible exit options'
])

add_title_slide(prs, 'Be a Partner in Success', 'Real Estate Asset Revitalization Company')

output_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\canva-assets\business-plan-presentation-english.pptx'
prs.save(output_path)
print('English PowerPoint created successfully')
