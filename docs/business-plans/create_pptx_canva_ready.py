from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

BONDS_GOLD = RGBColor(212, 168, 83)
BONDS_DARK = RGBColor(10, 15, 26)
BONDS_GRAY = RGBColor(148, 163, 184)
BONDS_BLUE = RGBColor(0, 51, 102)

AR_NAME = 'د. طلال بن حسن الزهراني'
AR_TITLE = 'المستشار المالي'
EN_NAME = 'Dr. Talal bin Hassan Al-Zahrani'
EN_TITLE = 'Financial Advisor'
PHONE = '+966 56 756 6616'
EMAIL = 'info@bonds-global.com'

def add_footer(slide, prs, lang='ar'):
    footer_box = slide.shapes.add_textbox(0, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.4))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    if lang == 'ar':
        p.text = f'{AR_TITLE} | {AR_NAME} | {PHONE} | {EMAIL}'
        p.alignment = PP_ALIGN.CENTER
    else:
        p.text = f'{EN_TITLE} | {EN_NAME} | {PHONE} | {EMAIL}'
        p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(9)
    p.font.color.rgb = BONDS_GRAY
    p.font.name = 'Arial'

def add_title_slide(prs, title, subtitle, lang='ar'):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = BONDS_DARK
    background.line.fill.background()
    
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(4.25), Inches(0.2), width=Inches(1.5))
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1.0))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = BONDS_GOLD
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(9), Inches(0.7))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    if lang == 'ar':
        lines = [
            ('المستشار المالي', 14, True, BONDS_GOLD),
            ('د. طلال بن حسن الزهراني', 18, True, RGBColor(255, 255, 255)),
            ('+966 56 756 6616', 12, False, BONDS_GRAY),
            ('info@bonds-global.com', 12, False, BONDS_GOLD),
        ]
    else:
        lines = [
            ('Financial Advisor', 14, True, BONDS_GOLD),
            ('Dr. Talal bin Hassan Al-Zahrani', 18, True, RGBColor(255, 255, 255)),
            ('+966 56 756 6616', 12, False, BONDS_GRAY),
            ('info@bonds-global.com', 12, False, BONDS_GOLD),
        ]
    
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(1.2))
    tf = contact_box.text_frame
    tf.word_wrap = True
    for i, (text, size, bold, color) in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    
    # QR code
    qr_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-qr-code.png'
    if os.path.exists(qr_path):
        slide.shapes.add_picture(qr_path, Inches(4.4), Inches(6.3), width=Inches(0.7))
    
    p = slide.shapes.add_textbox(Inches(0.5), Inches(7.0), Inches(9), Inches(0.3)).text_frame.paragraphs[0]
    p.text = 'www.bonds-global.com'
    p.font.size = Pt(9)
    p.font.color.rgb = BONDS_GRAY
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, lang='ar'):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = BONDS_DARK
    header.line.fill.background()
    
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.1), width=Inches(0.7))
    
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.2), Inches(8.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BONDS_GOLD
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.RIGHT if lang == 'ar' else PP_ALIGN.LEFT
    
    content_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.3), Inches(8.8), Inches(5.3))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = '• ' + bullet
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.RIGHT if lang == 'ar' else PP_ALIGN.LEFT
        p.space_after = Pt(10)
    
    add_footer(slide, prs, lang)
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items, lang='ar'):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.0))
    header.fill.solid()
    header.fill.fore_color.rgb = BONDS_DARK
    header.line.fill.background()
    
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.1), width=Inches(0.7))
    
    title_box = slide.shapes.add_textbox(Inches(1.1), Inches(0.2), Inches(8.4), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = BONDS_GOLD
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.RIGHT if lang == 'ar' else PP_ALIGN.LEFT
    
    align = PP_ALIGN.RIGHT if lang == 'ar' else PP_ALIGN.LEFT
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(4.3), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BONDS_BLUE
    p.font.name = 'Arial'
    p.alignment = align
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = '• ' + item
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = align
        p.space_after = Pt(6)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.3), Inches(4.3), Inches(5.2))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = BONDS_BLUE
    p.font.name = 'Arial'
    p.alignment = align
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = '• ' + item
        p.font.size = Pt(15)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = align
        p.space_after = Pt(6)
    
    add_footer(slide, prs, lang)
    return slide

def create_presentation(lang='ar'):
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    if lang == 'ar':
        add_title_slide(prs, 'خطة عمل', 'شركة إحياء الأصول غير المكتملة العقارية', lang='ar')
        
        add_content_slide(prs, 'ملخص تنفيذي', [
            'نموذج استثماري فريد لإحياء العقارات غير المكتملة',
            'شراكات محاصة مع المالكين دون شراء أراضٍ',
            'تمويل ذاتي من المساهمين المؤسسين',
            'الالتزام بالكود السعودي والتأمين الشامل',
            'عائد داخلي متوقع 19% - 20%'
        ], lang='ar')
        
        add_two_column_slide(prs, 'المشكلة والفرصة',
            'المشكلة', ['آلاف العقارات غير المكتملة', 'نقص السيولة', 'صعوبة الوصول للمقاولين'],
            'الفرصة', ['رؤية 2030', 'الطلب يفوق العرض', 'عوائد إيجارية مرتفعة'], lang='ar')
        
        add_content_slide(prs, 'الحل', [
            'البحث عن العقارات غير المكتملة',
            'التقييم الهندسي والقانوني',
            'عقد محاصة واضح',
            'التشطيب بأعلى معايير',
            'التسويق والبيع',
            'توزيع الأرباح'
        ], lang='ar')
        
        add_content_slide(prs, 'المؤشرات المالية', [
            'رأس المال المصرح: 200 مليون ريال',
            'الاستثمار المستهدف: 226 مليون ريال',
            'الاستثمار المبدئي: 25 مليون ريال',
            'IRR: 19% - 20%',
            'عائد الخروج: 2.5x - 4x'
        ], lang='ar')
        
        add_content_slide(prs, 'نموذج العمل', [
            'المحاصة الكاملة',
            'الشراء الجزئي',
            'رسوم الإدارة',
            'النموذج المختلط'
        ], lang='ar')
        
        add_content_slide(prs, 'خطة التنفيذ', [
            'السنة 1-2: إثبات النموذج',
            'السنة 3-5: النمو والتكرار',
            'السنة 5-7: الخروج الاستراتيجي'
        ], lang='ar')
        
        add_content_slide(prs, 'إدارة المخاطر', [
            'تأمين شامل',
            'فحص هندسي وقانوني',
            'عقود واضحة مع تحكيم ملزم',
            'كفالات بنكية',
            'احتياطي مخاطر 10%'
        ], lang='ar')
        
        add_title_slide(prs, 'كن شريكاً في النجاح', 'شركة إحياء الأصول غير المكتملة العقارية', lang='ar')
    else:
        add_title_slide(prs, 'Business Plan', 'Real Estate Asset Revitalization Company', lang='en')
        
        add_content_slide(prs, 'Executive Summary', [
            'Unique investment model for incomplete assets',
            'Partnership agreements without land purchase',
            'Self-funded by founding shareholders',
            'Saudi Building Code and insurance compliance',
            'Expected IRR 19% - 20%'
        ], lang='en')
        
        add_two_column_slide(prs, 'Problem & Opportunity',
            'Problem', ['Thousands of incomplete assets', 'Liquidity shortage', 'Reliable contractor access'],
            'Opportunity', ['Vision 2030', 'Demand exceeds supply', 'High rental yields'], lang='en')
        
        add_content_slide(prs, 'Solution', [
            'Identify incomplete assets',
            'Engineering and legal assessment',
            'Clear partnership contract',
            'High-standard finishing',
            'Marketing and sales',
            'Profit distribution'
        ], lang='en')
        
        add_content_slide(prs, 'Financial Highlights', [
            'Authorized Capital: SAR 200 million',
            'Target Investment: SAR 226 million',
            'Initial Investment: SAR 25 million',
            'IRR: 19% - 20%',
            'Exit Return: 2.5x - 4x'
        ], lang='en')
        
        add_content_slide(prs, 'Business Model', [
            'Full Muhassah',
            'Partial Purchase',
            'Management Fee',
            'Hybrid Model'
        ], lang='en')
        
        add_content_slide(prs, 'Implementation', [
            'Years 1-2: Proof of Concept',
            'Years 3-5: Growth and Replication',
            'Years 5-7: Strategic Exit'
        ], lang='en')
        
        add_content_slide(prs, 'Risk Management', [
            'Comprehensive insurance',
            'Engineering and legal inspection',
            'Clear contracts with arbitration',
            'Bank guarantees',
            '10% risk reserve'
        ], lang='en')
        
        add_title_slide(prs, 'Be a Partner in Success', 'Real Estate Asset Revitalization Company', lang='en')
    
    return prs

# Create both presentations
arabic_prs = create_presentation('ar')
arabic_prs.save(r'C:\Users\vip\bonds-global-web\docs\business-plans\canva-ready-arabic\business-plan-presentation-arabic.pptx')
print('Arabic Canva-ready presentation created')

english_prs = create_presentation('en')
english_prs.save(r'C:\Users\vip\bonds-global-web\docs\business-plans\canva-ready-english\business-plan-presentation-english.pptx')
print('English Canva-ready presentation created')
