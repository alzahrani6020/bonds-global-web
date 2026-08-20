from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
import numpy as np
from PIL import Image
import os

# Bonds brand colors
BONDS_GOLD = '#D4A853'
BONDS_DARK = '#0A0F1A'
BONDS_GRAY = '#94A3B8'
BONDS_BLUE = '#003366'
BONDS_LIGHT_GOLD = '#F0C96A'

output_dir = r'C:\Users\vip\bonds-global-web\docs\business-plans\canva-assets'
os.makedirs(output_dir, exist_ok=True)

def add_title_slide(prs, title, subtitle, lang='ar'):
    blank_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(blank_layout)
    
    # Background
    background = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(10, 15, 26)
    background.line.fill.background()
    
    # Logo
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(4.25), Inches(0.3), width=Inches(1.5))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 168, 83)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.7), Inches(9), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    
    # Contact info
    if lang == 'ar':
        contact_lines = ['المستشار المالي', 'د. طلال بن حسن الزهراني', '+966 56 756 6616']
    else:
        contact_lines = ['Financial Advisor', 'Dr. Talal bin Hassan Al-Zahrani', '+966 56 756 6616']
    
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.8), Inches(9), Inches(1.2))
    tf = contact_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(contact_lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        if i == 0:
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = RGBColor(212, 168, 83)
        elif i == 1:
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
        else:
            p.font.size = Pt(12)
            p.font.color.rgb = RGBColor(148, 163, 184)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets, layout='content'):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Header bar
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(10, 15, 26)
    header.line.fill.background()
    
    # Logo small
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.15), width=Inches(0.8))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.25), Inches(8.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 168, 83)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.RIGHT
    
    # Content
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(8.6), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = '• ' + bullet
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Header bar
    header = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(10, 15, 26)
    header.line.fill.background()
    
    logo_path = r'C:\Users\vip\bonds-global-web\docs\business-plans\bonds-logo.png'
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.3), Inches(0.15), width=Inches(0.8))
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1.2), Inches(0.25), Inches(8.3), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(212, 168, 83)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.RIGHT
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.RIGHT
    
    for item in left_items:
        p = tf.add_paragraph()
        p.text = '• ' + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(4.2), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 51, 102)
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.RIGHT
    
    for item in right_items:
        p = tf.add_paragraph()
        p.text = '• ' + item
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(10, 15, 26)
        p.font.name = 'Arial'
        p.alignment = PP_ALIGN.RIGHT
        p.space_after = Pt(8)
    
    return slide

def create_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, 'خطة عمل', 'شركة إحياء الأصول غير المكتملة العقارية', lang='ar')
    
    # Slide 2: Executive Summary
    add_content_slide(prs, 'ملخص تنفيذي', [
        'نموذج استثماري فريد لإحياء العقارات غير المكتملة في المملكة',
        'شراكات محاصة مع المالكين دون شراء أراضٍ',
        'تمويل ذاتي من المساهمين المؤسسين',
        'الالتزام بالكود السعودي والتأمين الشامل',
        'عائد داخلي متوقع 19% - 20% على أفق 10 سنوات'
    ])
    
    # Slide 3: Problem & Opportunity
    add_two_column_slide(prs, 'المشكلة والفرصة', 
        'المشكلة', [
            'آلاف العقارات غير المكتملة',
            'نقص السيولة لدى المالكين',
            'صعوبة الوصول للمقاولين الموثوقين',
            'تعقيد الإجراءات القانونية'
        ],
        'الفرصة', [
            'رؤية 2030 ونمو القطاع العقاري',
            'الطلب يفوق العرض في المدن الرئيسية',
            'عوائد إيجارية مرتفعة',
            'نموذج محاصة مرن ومربح'
        ])
    
    # Slide 4: Solution
    add_content_slide(prs, 'الحل: نموذج العمل', [
        'نبحث عن العقارات غير المكتملة في المواقع المستهدفة',
        'نقيّم العقار هندسياً وقانونياً ومالياً',
        'نُبرم عقد محاصة واضح مع المالك',
        'نُنجز التشطيب وفق أعلى معايير الجودة',
        'نُسوّق الوحدات ونبيعها بأسعار السوق',
        'نوزع الأرباح بين الشركة والمالك'
    ])
    
    # Slide 5: Market Indicators
    add_content_slide(prs, 'مؤشرات السوق العقاري', [
        'عائد إيجاري في الرياض: 8.89%',
        'عائد إيجاري في جدة: 7.89%',
        'نمو الصفقات العقارية (2025): +14.2%',
        'حجم التمويل العقاري: 730 مليار ريال',
        'فجوة السعر الفعلي/الطلب: 6%',
        'علاوة البناء الجديد: +25%'
    ])
    
    # Slide 6: Target Areas
    add_two_column_slide(prs, 'المناطق المستهدفة',
        'الرياض', [
            'الشمال الفاخر: العقيق، حطين، الملقا',
            'الشمال الشرقي: المونسية، الرمال',
            'الياسمين، النرجس',
            'عوائد إيجارية مرتفعة'
        ],
        'جدة', [
            'الساحل الشمالي: أبحر، الصواري',
            'الوسط الحيوي: الحمراء، النعيم',
            'الزهراء، السلامة',
            'طلب سكني وتجاري مستمر'
        ])
    
    # Slide 7: Partnership Models
    add_content_slide(prs, 'نماذج الشراكة', [
        'المحاصة الكاملة: نتولى التشطيب والتسويق مقابل نسبة من الأرباح',
        'الشراء الجزئي: نشتري جزءاً من الوحدات بعد التشطيب',
        'رسوم الإدارة: رسوم ثابتة مقابل التشطيب والتسويق',
        'النموذج المختلط: مزيج من المحاصة والشراء ورسوم الإدارة'
    ])
    
    # Slide 8: Financial Highlights
    add_content_slide(prs, 'المؤشرات المالية الرئيسية', [
        'رأس المال المصرح: 200 مليون ريال',
        'إجمالي الاستثمار المستهدف: 226 مليون ريال',
        'الاستثمار المبدئي: 25 مليون ريال',
        'IRR المتوقع: 19% - 20%',
        'مضاعفة رأس المال عند الخروج: 2.5x - 4x',
        'فترة الاسترداد: 5.5 سنة'
    ])
    
    # Slide 9: Revenue Model
    add_content_slide(prs, 'مصادر الإيرادات', [
        'إيرادات المحاصة (نصيب الشركة): 60%',
        'أرباح إعادة البيع والشراء: 25%',
        'رسوم الإدارة والتسويق: 10%',
        'خدمات استشارية وإدارية: 5%'
    ])
    
    # Slide 10: Implementation Timeline
    add_content_slide(prs, 'المراحل الزمنية', [
        'السنة 1 - 2: إثبات النموذج (3 - 5 مشاريع)',
        'السنة 3 - 5: النمو والتكرار (6 - 10 مشاريع سنوياً)',
        'السنة 5 - 7: الخروج الاستراتيجي',
        'الخيارات: بيع حصة مسيطرة، طرح جزئي، أو إعادة شراء'
    ])
    
    # Slide 11: Risk Mitigation
    add_content_slide(prs, 'إدارة المخاطر', [
        'تأمين شامل (CAR، مسؤولية مدنية، عيوب خفية)',
        'فحص هندسي وقانوني قبل التعاقد',
        'عقود واضحة مع شرط تحكيم ملزم',
        'كفالات بنكية من المقاولين',
        'احتياطي مخاطر 10%',
        'تنويع جغرافي وفئات استهداف'
    ])
    
    # Slide 12: Team
    add_content_slide(prs, 'الفريق الإداري', [
        'الرئيس التنفيذي: الاستراتيجية وعلاقات المستثمرين',
        'المدير التشغيلي: إدارة المشاريع والجودة',
        'المدير المالي: التدفقات النقدية والتقارير',
        'المدير التجاري: التسويق والمبيعات',
        'المدير القانوني: العقود والامتثال',
        'خطة التوظيف: 10 - 15 موظفاً في السنة الأولى'
    ])
    
    # Slide 13: Why Invest
    add_content_slide(prs, 'لماذا نستثمر؟', [
        'نموذج فريد لا يتطلب شراء أرض',
        'سوق ضخم من الأصول غير المكتملة',
        'تمويل ذاتي بدون ديون أو صكوك',
        'عوائد جذابة ومخاطر محسوبة',
        'فريق إداري متخصص وشبكة علاقات قوية',
        'خيارات خروج واضحة ومرنة'
    ])
    
    # Slide 14: Call to Action
    add_title_slide(prs, 'كن شريكاً في النجاح', 'شركة إحياء الأصول غير المكتملة العقارية', lang='ar')
    
    prs.save(os.path.join(output_dir, 'business-plan-presentation.pptx'))
    print('PowerPoint presentation created')

def create_slide_content_md():
    content = """# محتوى شرائح العرض التقديمي لـ Canva

## الشريحة 1: الغلاف
- **العنوان:** خطة عمل
- **العنوان الفرعي:** شركة إحياء الأصول غير المكتملة العقارية
- **الشعار:** شعار بوندز
- **الألوان:** خلفية داكنة #0A0F1A، نص ذهبي #D4A853

## الشريحة 2: ملخص تنفيذي
- نموذج استثماري فريد لإحياء العقارات غير المكتملة في المملكة
- شراكات محاصة مع المالكين دون شراء أراضٍ
- تمويل ذاتي من المساهمين المؤسسين
- الالتزام بالكود السعودي والتأمين الشامل
- عائد داخلي متوقع 19% - 20% على أفق 10 سنوات

## الشريحة 3: المشكلة والفرصة
**المشكلة:**
- آلاف العقارات غير المكتملة
- نقص السيولة لدى المالكين
- صعوبة الوصول للمقاولين الموثوقين

**الفرصة:**
- رؤية 2030 ونمو القطاع العقاري
- الطلب يفوق العرض
- عوائد إيجارية مرتفعة

## الشريحة 4: الحل
- البحث عن العقارات غير المكتملة
- التقييم الهندسي والقانوني والمالي
- عقد محاصة واضح
- التشطيب بأعلى معايير الجودة
- التسويق والبيع
- توزيع الأرباح

## الشريحة 5: مؤشرات السوق
- عائد إيجاري في الرياض: 8.89%
- عائد إيجاري في جدة: 7.89%
- نمو الصفقات العقارية: +14.2%
- حجم التمويل العقاري: 730 مليار ريال

## الشريحة 6: المناطق المستهدفة
**الرياض:** الشمال الفاخر، الشمال الشرقي
**جدة:** الساحل الشمالي، الوسط الحيوي

## الشريحة 7: نماذج الشراكة
- المحاصة الكاملة
- الشراء الجزئي
- رسوم الإدارة
- النموذج المختلط

## الشريحة 8: المؤشرات المالية
- رأس المال المصرح: 200 مليون ريال
- الاستثمار المستهدف: 226 مليون ريال
- IRR: 19% - 20%
- مضاعفة رأس المال: 2.5x - 4x

## الشريحة 9: مصادر الإيرادات
- المحاصة: 60%
- إعادة البيع: 25%
- رسوم الإدارة: 10%
- خدمات استشارية: 5%

## الشريحة 10: المراحل الزمنية
- السنة 1-2: إثبات النموذج
- السنة 3-5: النمو والتكرار
- السنة 5-7: الخروج الاستراتيجي

## الشريحة 11: إدارة المخاطر
- تأمين شامل
- فحص هندسي وقانوني
- عقود واضحة مع تحكيم ملزم
- كفالات بنكية

## الشريحة 12: الفريق الإداري
- الرئيس التنفيذي
- المدير التشغيلي
- المدير المالي
- المدير التجاري
- المدير القانوني

## الشريحة 13: لماذا نستثمر؟
- نموذج فريد
- سوق ضخم
- تمويل ذاتي
- عوائد جذابة

## الشريحة 14: الدعوة للاستثمار
- **العنوان:** كن شريكاً في النجاح
- **النص:** شركة إحياء الأصول غير المكتملة العقارية
- **زر:** تواصل معنا
"""
    
    with open(os.path.join(output_dir, 'slide-content-for-canva.md'), 'w', encoding='utf-8') as f:
        f.write(content)
    print('Slide content markdown created')

def create_charts():
    plt.rcParams['font.family'] = 'Arial'
    
    # Chart 1: Revenue Projection
    years = ['Year 1', 'Year 2', 'Year 3', 'Year 4', 'Year 5']
    revenues = [3.2, 6.4, 9.6, 12.8, 16.0]
    profits = [-0.4, 2.35, 4.79, 7.23, 9.77]
    
    fig, ax = plt.subplots(figsize=(10, 6))
    x = np.arange(len(years))
    width = 0.35
    bars1 = ax.bar(x - width/2, revenues, width, label='Revenue', color=BONDS_GOLD)
    bars2 = ax.bar(x + width/2, profits, width, label='Net Profit', color=BONDS_DARK)
    ax.set_xlabel('Year', fontsize=12)
    ax.set_ylabel('SAR Million', fontsize=12)
    ax.set_title('Financial Projections (SAR Million)', fontsize=16, fontweight='bold', color=BONDS_DARK)
    ax.set_xticks(x)
    ax.set_xticklabels(years)
    ax.legend()
    ax.grid(axis='y', alpha=0.3)
    plt.tight_layout()
    plt.savefig(os.path.join(output_dir, 'chart-revenue-projection.png'), dpi=300, bbox_inches='tight')
    plt.close()
    
    # Chart 2: Revenue Sources Pie
    fig, ax = plt.subplots(figsize=(8, 8))
    labels = ['Muhassah\n60%', 'Resale\n25%', 'Management Fees\n10%', 'Consulting\n5%']
    sizes = [60, 25, 10, 5]
    colors = [BONDS_GOLD, BONDS_DARK, BONDS_LIGHT_GOLD, BONDS_GRAY]
    explode = (0.05, 0, 0, 0)
    ax.pie(sizes, explode=explode, labels=labels, colors=colors, autopct='%1.0f%%',
           shadow=False, startangle=90, textprops={'fontsize': 12})
    ax.set_title('Revenue Sources', fontsize=16, fontweight='bold', color=BONDS_DARK)
    plt.tight_layout()
    plt.savefig(os.path.join(output_dir, 'chart-revenue-sources.png'), dpi=300, bbox_inches='tight')
    plt.close()
    
    # Chart 3: Capital Utilization
    fig, ax = plt.subplots(figsize=(8, 8))
    labels = ['Project Investments\n66%', 'Working Capital\n18%', 'Establishment\n7%', 'Technology\n4%', 'Reserve\n5%']
    sizes = [66, 18, 7, 4, 5]
    colors = [BONDS_DARK, BONDS_GOLD, BONDS_LIGHT_GOLD, BONDS_GRAY, '#5C6B7F']
    ax.pie(sizes, labels=labels, colors=colors, autopct='%1.0f%%',
           shadow=False, startangle=90, textprops={'fontsize': 11})
    ax.set_title('Capital Utilization', fontsize=16, fontweight='bold', color=BONDS_DARK)
    plt.tight_layout()
    plt.savefig(os.path.join(output_dir, 'chart-capital-utilization.png'), dpi=300, bbox_inches='tight')
    plt.close()
    
    # Chart 4: Implementation Timeline
    fig, ax = plt.subplots(figsize=(12, 5))
    phases = ['Establishment', 'Launch', 'Growth', 'Expansion', 'Exit']
    starts = [0, 0.5, 2, 4, 5]
    durations = [0.5, 1, 2, 2, 2]
    colors = [BONDS_DARK, BONDS_GOLD, BONDS_LIGHT_GOLD, BONDS_GRAY, '#5C6B7F']
    
    for i, (phase, start, duration, color) in enumerate(zip(phases, starts, durations, colors)):
        ax.barh(i, duration, left=start, height=0.5, color=color, edgecolor='white')
        ax.text(start + duration/2, i, phase, ha='center', va='center', color='white', fontweight='bold', fontsize=10)
    
    ax.set_yticks(range(len(phases)))
    ax.set_yticklabels([''] * len(phases))
    ax.set_xlabel('Years', fontsize=12)
    ax.set_title('Implementation Timeline', fontsize=16, fontweight='bold', color=BONDS_DARK)
    ax.set_xlim(0, 8)
    ax.grid(axis='x', alpha=0.3)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_visible(False)
    plt.tight_layout()
    plt.savefig(os.path.join(output_dir, 'chart-timeline.png'), dpi=300, bbox_inches='tight')
    plt.close()
    
    # Chart 5: Risk Matrix
    fig, ax = plt.subplots(figsize=(10, 7))
    risks = ['Asset Acquisition', 'Hidden Defects', 'Sales Delay', 'Contractor Delay', 'Material Prices', 'Legal Disputes', 'Price Decline']
    probability = [0.9, 0.9, 0.5, 0.5, 0.5, 0.5, 0.3]
    impact = [0.8, 0.9, 0.5, 0.4, 0.4, 0.6, 0.4]
    colors = ['red' if p > 0.7 or i > 0.7 else 'orange' if p > 0.4 or i > 0.4 else 'green' for p, i in zip(probability, impact)]
    
    scatter = ax.scatter(probability, impact, s=300, c=colors, alpha=0.7, edgecolors='black')
    for i, risk in enumerate(risks):
        ax.annotate(risk, (probability[i], impact[i]), xytext=(5, 5), textcoords='offset points', fontsize=9)
    
    ax.set_xlabel('Probability', fontsize=12)
    ax.set_ylabel('Impact', fontsize=12)
    ax.set_title('Risk Matrix', fontsize=16, fontweight='bold', color=BONDS_DARK)
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.grid(True, alpha=0.3)
    ax.axhline(y=0.5, color='gray', linestyle='--', alpha=0.5)
    ax.axvline(x=0.5, color='gray', linestyle='--', alpha=0.5)
    plt.tight_layout()
    plt.savefig(os.path.join(output_dir, 'chart-risk-matrix.png'), dpi=300, bbox_inches='tight')
    plt.close()
    
    # Chart 6: Market Indicators
    fig, ax = plt.subplots(figsize=(9, 6))
    indicators = ['Riyadh\nYield', 'Jeddah\nYield', 'Transaction\nGrowth', 'Price Gap', 'New Build\nPremium']
    values = [8.89, 7.89, 14.2, 6, 25]
    colors = [BONDS_GOLD, BONDS_GOLD, BONDS_DARK, BONDS_LIGHT_GOLD, BONDS_GRAY]
    bars = ax.bar(indicators, values, color=colors, edgecolor='white')
    ax.set_ylabel('Percentage (%)', fontsize=12)
    ax.set_title('Real Estate Market Indicators', fontsize=16, fontweight='bold', color=BONDS_DARK)
    ax.grid(axis='y', alpha=0.3)
    for bar in bars:
        height = bar.get_height()
        ax.annotate(f'{height}%', xy=(bar.get_x() + bar.get_width() / 2, height),
                    xytext=(0, 3), textcoords="offset points", ha='center', va='bottom', fontweight='bold')
    plt.tight_layout()
    plt.savefig(os.path.join(output_dir, 'chart-market-indicators.png'), dpi=300, bbox_inches='tight')
    plt.close()
    
    print('Charts created')

def create_design_guidelines():
    content = """# دليل تصميم هوية بوندز للعرض التقديمي

## ألوان العلامة التجارية

| اللون | الكود | الاستخدام |
|-------|-------|-----------|
| الذهبي | #D4A853 | العناوين الرئيسية، الأزرار، العناصر البارزة |
| الذهبي الفاتح | #F0C96A | التأكيدات والتفاصيل الثانوية |
| الأزرق الداكن | #0A0F1A | الخلفيات الداكنة، النصوص الرئيسية |
| الأزرق العميق | #003366 | عناوين فرعية |
| الرمادي | #94A3B8 | النصوص الفرعية والتوضيحية |

## الخطوط المقترحة

- **العربية:** Vazirmatn, Cairo, Arial
- **الإنجليزية:** Inter, Arial, Montserrat

## مقاسات الشرائح

- **عرض تقديمي:** 1920 × 1080 بكسل (16:9)
- **بوست إنستغرام:** 1080 × 1080 بكسل (1:1)
- **ستوري:** 1080 × 1920 بكسل (9:16)
- **A4:** 2480 × 3508 بكسل

## عناصر التصميم

### الترويسة
- شعار بوندز في الزاوية العلوية
- خط فاصل ذهبي سمكه 2 بكسل
- خلفية داكنة للترويسة

### التذييل
- Bonds Financial Management & Consulting
- www.bonds-global.com
- خط ذهبي فاصل

### الجداول
- رأس جدول بخلفية داكنة #0A0F1A
- نص رأس الجدول باللون الأبيض
- صفوف متناوبة باللونين الأبيض والذهبي الفاتح جداً #FDF8F0

### الأيقونات
- أيقونات بسيطة ورفيعة
- باللون الذهبي أو الأزرق الداكن
- متناسقة مع الطابع الاحترافي

## قواعد استخدام الشعار

- لا تشوه نسب الشعار
- اترك مساحة فارغة حول الشعار
- لا تستخدم الشعار على خلفيات معقدة
- الشعار يعمل بشكل أفضل على الخلفيات الداكنة

## نمط الصور

- صور عقارات عالية الجودة
- إضاءة طبيعية
- تركيز على "قبل وبعد" للمشاريع
- تعديل الألوان لتكون دافئة وذهبية

## نصائح للعرض التقديمي

1. لا تزيد عن 6 نقاط في الشريحة الواحدة
2. استخدم الصور الكبيرة والواضحة
3. اترك مساحة فارغة كافية
4. استخدم الرسوم البيانية لتبسيط الأرقام
5. حافظ على تناسق الألوان في جميع الشرائح
"""
    
    with open(os.path.join(output_dir, 'design-guidelines.md'), 'w', encoding='utf-8') as f:
        f.write(content)
    print('Design guidelines created')

# Execute all
create_pptx()
create_slide_content_md()
create_charts()
create_design_guidelines()
print('All Canva assets created successfully')
